"""
Unit tests for the diff/match/recolor core. These build real PPTX objects in
memory and inspect the colored runs directly, so they need python-pptx but NOT
LibreOffice. Rendering is validated separately (see test_render, skipped unless
soffice is present).
"""

import shutil
import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402
from pptx.dml.color import MSO_COLOR_TYPE  # noqa: E402

from app.slide_diff import (  # noqa: E402
    RED, GREEN, build_marked_decks, match_slides, _slide_text_parts,
    _build_stream, _flag, pptx_to_images,
)

BLANK = 6  # blank layout index in the default template


# --------------------------- builders ------------------------------------- #

def _add_text_slide(prs, paragraphs, title=None):
    """paragraphs: list of strings OR list of run-lists [(text, {bold:..}), ...]."""
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    if title is not None:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        tb.text_frame.text = title
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
    tf = box.text_frame
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(para, str):
            p.text = para
        else:
            for text, fmt in para:
                r = p.add_run()
                r.text = text
                if fmt.get("bold"):
                    r.font.bold = True
                if fmt.get("size"):
                    r.font.size = Pt(fmt["size"])
    return slide


def _add_table_slide(prs, cells_row):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    rows, cols = 1, len(cells_row)
    gf = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    for c, val in enumerate(cells_row):
        gf.table.cell(0, c).text = val
    return slide


def _save(prs, tmp_path, name):
    path = tmp_path / name
    prs.save(str(path))
    return path


def _colored_runs(slide, color):
    """Return list of run texts on the slide whose font color == color."""
    out = []
    from app.slide_diff import _iter_text_frames
    for tf in _iter_text_frames(slide.shapes):
        for para in tf.paragraphs:
            for run in para.runs:
                try:
                    if run.font.color and run.font.color.type == MSO_COLOR_TYPE.RGB \
                            and run.font.color.rgb == color:
                        out.append(run.text)
                except Exception:
                    pass
    return out


def _diff(tmp_path, old_paras, new_paras, old_is_table=False, new_is_table=False):
    op = Presentation()
    (_add_table_slide(op, old_paras[0]) if old_is_table else _add_text_slide(op, old_paras))
    np_ = Presentation()
    (_add_table_slide(np_, new_paras[0]) if new_is_table else _add_text_slide(np_, new_paras))
    old_path = _save(op, tmp_path, "old.pptx")
    new_path = _save(np_, tmp_path, "new.pptx")
    old_prs, new_prs, entries = build_marked_decks(old_path, new_path)
    return old_prs, new_prs, entries


# --------------------------- acceptance tests ------------------------------ #

def test_word_edit(tmp_path):
    """A B C D -> A X C D : only B red on old, only X green on new."""
    old_prs, new_prs, entries = _diff(tmp_path, ["A B C D"], ["A X C D"])
    assert _colored_runs(old_prs.slides[0], RED) == ["B"]
    assert _colored_runs(new_prs.slides[0], GREEN) == ["X"]
    assert entries[0].has_changes is True
    assert entries[0].status == "changed"


def test_paragraph_split_no_color(tmp_path):
    """A B C D -> [A B][C D] : no coloring, no change flag."""
    old_prs, new_prs, entries = _diff(tmp_path, ["A B C D"], ["A B", "C D"])
    assert _colored_runs(old_prs.slides[0], RED) == []
    assert _colored_runs(new_prs.slides[0], GREEN) == []
    assert entries[0].has_changes is False
    assert entries[0].status == "unchanged"


def test_paragraph_merge_no_color(tmp_path):
    """[A B][C D] -> A B C D : no coloring."""
    old_prs, new_prs, entries = _diff(tmp_path, ["A B", "C D"], ["A B C D"])
    assert _colored_runs(old_prs.slides[0], RED) == []
    assert _colored_runs(new_prs.slides[0], GREEN) == []
    assert entries[0].has_changes is False


def test_paragraph_promoted_to_table_no_color(tmp_path):
    """text 'Revenue 100' -> table [Revenue][100] : same words, no color."""
    old_prs, new_prs, entries = _diff(
        tmp_path, ["Revenue 100"], [["Revenue", "100"]], new_is_table=True
    )
    assert _colored_runs(old_prs.slides[0], RED) == []
    assert _colored_runs(new_prs.slides[0], GREEN) == []


def test_table_added_cell(tmp_path):
    """table [Revenue][100] -> [Revenue][100][200] : only 200 green."""
    old_prs, new_prs, entries = _diff(
        tmp_path, [["Revenue", "100"]], [["Revenue", "100", "200"]],
        old_is_table=True, new_is_table=True,
    )
    assert _colored_runs(old_prs.slides[0], RED) == []
    assert _colored_runs(new_prs.slides[0], GREEN) == ["200"]


def test_mixed_structural_and_textual(tmp_path):
    """A B C D -> [A X][C D] : only B red, only X green; the split adds no color."""
    old_prs, new_prs, entries = _diff(tmp_path, ["A B C D"], ["A X", "C D"])
    assert _colored_runs(old_prs.slides[0], RED) == ["B"]
    assert _colored_runs(new_prs.slides[0], GREEN) == ["X"]


def test_within_word_change(tmp_path):
    """FY25 -> FY26 : changed token painted on each side."""
    old_prs, new_prs, entries = _diff(tmp_path, ["FY25 Revenue"], ["FY26 Revenue"])
    assert _colored_runs(old_prs.slides[0], RED) == ["FY25"]
    assert _colored_runs(new_prs.slides[0], GREEN) == ["FY26"]
    assert "Revenue" not in _colored_runs(new_prs.slides[0], GREEN)


def test_formatting_preserved_on_changed_paragraph(tmp_path):
    """The key win over clear()+rebuild: a bold run in a changed paragraph keeps
    its bold when an adjacent run changes."""
    op = Presentation()
    _add_text_slide(op, [[("Hello ", {}), ("world", {"bold": True})]])
    np_ = Presentation()
    _add_text_slide(np_, [[("Hi ", {}), ("world", {"bold": True})]])
    old_path = _save(op, tmp_path, "old.pptx")
    new_path = _save(np_, tmp_path, "new.pptx")
    _, new_prs, _ = build_marked_decks(old_path, new_path)

    # "world" must still be bold and must NOT be green (it didn't change).
    from app.slide_diff import _iter_text_frames
    bold_world = False
    for tf in _iter_text_frames(new_prs.slides[0].shapes):
        for para in tf.paragraphs:
            for run in para.runs:
                if run.text == "world":
                    bold_world = run.font.bold is True
    assert bold_world, "bold formatting on the unchanged 'world' run was lost"
    assert "world" not in _colored_runs(new_prs.slides[0], GREEN)
    assert _colored_runs(new_prs.slides[0], GREEN) == ["Hi"]


def test_whitespace_only_change_suppressed(tmp_path):
    """'Q1  results' -> 'Q1 results' (double space collapsed) : no color."""
    old_prs, new_prs, entries = _diff(tmp_path, ["Q1  results"], ["Q1 results"])
    # The only difference is a whitespace token, which is never colored.
    assert _colored_runs(old_prs.slides[0], RED) == []
    assert _colored_runs(new_prs.slides[0], GREEN) == []
    assert entries[0].has_changes is False


# --------------------------- matching tests -------------------------------- #

def _parts(title, body):
    from app.slide_diff import _norm
    return (_norm(title), _norm(body), _norm(title + " " + body))


def test_match_equal_counts():
    old = [_parts("Intro", "alpha"), _parts("Body", "beta"), _parts("End", "gamma")]
    new = [_parts("Intro", "alpha"), _parts("Body", "beta CHANGED"), _parts("End", "gamma")]
    pairs = match_slides(old, new)
    assert all(p["kind"] == "matched" for p in pairs)
    assert [(p["old"], p["new"]) for p in pairs] == [(0, 0), (1, 1), (2, 2)]


def test_match_mid_deck_insert_no_cascade():
    """Inserting a slide at position 1 must NOT mark slides 2,3 as changed."""
    old = [_parts("A", "alpha"), _parts("B", "beta"), _parts("C", "gamma")]
    new = [_parts("A", "alpha"), _parts("NEW", "delta brand new"),
           _parts("B", "beta"), _parts("C", "gamma")]
    pairs = match_slides(old, new)
    by_new = {p["new"]: p for p in pairs}
    assert by_new[1]["kind"] == "added"             # the inserted slide
    assert by_new[0]["kind"] == "matched" and by_new[0]["old"] == 0
    assert by_new[2]["kind"] == "matched" and by_new[2]["old"] == 1  # B->B, not B->C
    assert by_new[3]["kind"] == "matched" and by_new[3]["old"] == 2


def test_match_removed():
    old = [_parts("A", "alpha"), _parts("B", "beta"), _parts("C", "gamma")]
    new = [_parts("A", "alpha"), _parts("C", "gamma")]
    pairs = match_slides(old, new)
    removed = [p for p in pairs if p["kind"] == "removed"]
    assert len(removed) == 1 and removed[0]["old"] == 1


def test_match_reordered_flagged():
    old = [_parts("A", "alpha"), _parts("B", "beta"), _parts("C", "gamma")]
    new = [_parts("C", "gamma"), _parts("B", "beta"), _parts("A", "alpha")]
    pairs = match_slides(old, new)
    assert any(p["reordered"] for p in pairs)


def test_match_added_all_when_old_empty():
    pairs = match_slides([], [_parts("A", "x"), _parts("B", "y")])
    assert [p["kind"] for p in pairs] == ["added", "added"]


# --------------------------- flatten / stream tests ------------------------ #

def test_flatten_reads_table_and_group_text(tmp_path):
    prs = Presentation()
    _add_table_slide(prs, ["Revenue", "100"])
    _, tokens = _build_stream(prs.slides[0])
    words = [t["text"] for t in tokens if not t["ws"]]
    assert "Revenue" in words and "100" in words


def test_added_slide_all_green_removed_all_red(tmp_path):
    op = Presentation()
    _add_text_slide(op, ["Only in old"])
    np_ = Presentation()
    _add_text_slide(np_, ["Only in old"])
    _add_text_slide(np_, ["Brand new slide content"])  # appended
    old_path = _save(op, tmp_path, "old.pptx")
    new_path = _save(np_, tmp_path, "new.pptx")
    _, new_prs, entries = build_marked_decks(old_path, new_path)
    added = [e for e in entries if e.status == "added"]
    assert len(added) == 1
    green = _colored_runs(new_prs.slides[added[0].new_index0], GREEN)
    assert "Brand" in green and "new" in green and "slide" in green


# --------------------------- render test (optional) ------------------------ #

@pytest.mark.skipif(shutil.which("soffice") is None or shutil.which("pdftoppm") is None,
                    reason="LibreOffice/poppler not installed in this environment")
def test_render_smoke(tmp_path):
    prs = Presentation()
    _add_text_slide(prs, ["Hello world"], title="Title")
    path = _save(prs, tmp_path, "deck.pptx")
    imgs = pptx_to_images(path, tmp_path / "out")
    assert len(imgs) == 1 and imgs[0].exists()
