"""
slide_diff.py — text-only PPTX comparison with in-place color marking.

Pipeline per slide pair:
  1. Flatten each slide (text frames + table cells + groups, in document order)
     into one ordered word-token stream. Paragraph boundaries are marked with a
     sentinel so the diff never merges words across a break.
  2. Diff the two token streams with difflib.SequenceMatcher (word granularity).
     This mechanical diff is the ONLY authority for what gets colored.
  3. Recolor the exact changed runs *in place* in the PPTX XML (deleted -> red on
     the old deck, inserted -> green on the new deck), splitting a run only at a
     change boundary and preserving every run's own formatting.

Slide pairing is done with an optimal assignment (Hungarian) over title/body
text similarity, so inserting or reordering slides does not cascade false diffs.

Rendering (pptx -> pdf -> png) is delegated to LibreOffice headless + pdftoppm
and is used ONLY to produce preview images; it never detects changes.

No LLM is involved in detection or coloring. See enrich.py for the optional,
advisory-only summary layer.
"""

from __future__ import annotations

import copy
import bisect
import re
import subprocess
import uuid
from dataclasses import dataclass
from difflib import SequenceMatcher
from pathlib import Path
from typing import Optional

import numpy as np
from scipy.optimize import linear_sum_assignment

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.text.text import _Run

# Slightly-off shades so a later color-based locator (if ever added) never
# collides with text that is genuinely pure red/green in the source deck.
RED = RGBColor(0xC0, 0x00, 0x01)
GREEN = RGBColor(0x00, 0x70, 0x01)

SEP = "\x00"  # paragraph-boundary sentinel; never colored, never written back

A_R = qn("a:r")
A_BR = qn("a:br")
A_FLD = qn("a:fld")
A_T = qn("a:t")
XML_SPACE = "{http://www.w3.org/XML/1998/namespace}space"

_WORD_RE = re.compile(r"\s+|\S+")


# --------------------------------------------------------------------------- #
# Flatten a slide into an ordered token stream
# --------------------------------------------------------------------------- #

def _iter_text_frames(shapes):
    """Yield every text frame on a slide, recursing into groups and tables,
    in document order. Uses getattr guards so shapes lacking has_table /
    has_text_frame (pictures, connectors, ...) are simply skipped."""
    for sh in shapes:
        try:
            is_group = sh.shape_type == MSO_SHAPE_TYPE.GROUP
        except Exception:
            is_group = False
        if is_group:
            yield from _iter_text_frames(sh.shapes)
        elif getattr(sh, "has_table", False):
            seen = set()  # merged cells repeat during iteration; dedupe by identity
            for row in sh.table.rows:
                for cell in row.cells:
                    cid = id(cell._tc)
                    if cid in seen:
                        continue
                    seen.add(cid)
                    yield cell.text_frame
        elif getattr(sh, "has_text_frame", False):
            yield sh.text_frame


def _para_text_and_runs(para):
    """Return (text, run_map) for a paragraph, built by walking content children
    directly so the character offsets stay exactly aligned with the runs we will
    later color.

    run_map entries: (start, end, element, tag) over the paragraph's text.
    Only tag == A_R is colorable. a:br -> "\\n", a:fld text is read but locked.
    """
    p = para._p
    parts: list[str] = []
    run_map: list[tuple[int, int, object, str]] = []
    cursor = 0
    for child in p.iterchildren():
        tag = child.tag
        if tag == A_R or tag == A_FLD:
            t_el = child.find(A_T)
            t = t_el.text if (t_el is not None and t_el.text is not None) else ""
        elif tag == A_BR:
            t = "\n"
        else:
            continue  # a:pPr, a:endParaRPr, ... contribute no text
        parts.append(t)
        run_map.append((cursor, cursor + len(t), child, tag))
        cursor += len(t)
    return "".join(parts), run_map


@dataclass
class _ParaInfo:
    para: object
    text: str
    run_map: list


def _build_stream(slide):
    """Return (para_infos, tokens).
    token = {"text", "pidx", "ws"}; pidx is None for the SEP sentinel."""
    para_infos: list[_ParaInfo] = []
    tokens: list[dict] = []
    if slide is None:
        return para_infos, tokens
    for tf in _iter_text_frames(slide.shapes):
        for para in tf.paragraphs:
            text, run_map = _para_text_and_runs(para)
            pidx = len(para_infos)
            para_infos.append(_ParaInfo(para=para, text=text, run_map=run_map))
            for piece in _WORD_RE.findall(text):
                tokens.append({"text": piece, "pidx": pidx, "ws": piece.isspace()})
            tokens.append({"text": SEP, "pidx": None, "ws": True})
    return para_infos, tokens


def _stream_plaintext(tokens) -> str:
    """Human-readable slide text (sentinels dropped, breaks normalized)."""
    return "".join(t["text"] for t in tokens if t["pidx"] is not None)


# --------------------------------------------------------------------------- #
# Diff two token streams into per-token flags
# --------------------------------------------------------------------------- #

def _flag(old_tokens, new_tokens):
    a = [t["text"] for t in old_tokens]
    b = [t["text"] for t in new_tokens]
    old_del = [False] * len(old_tokens)
    new_ins = [False] * len(new_tokens)
    for tag, i1, i2, j1, j2 in SequenceMatcher(a=a, b=b, autojunk=False).get_opcodes():
        if tag in ("delete", "replace"):
            for i in range(i1, i2):
                old_del[i] = True
        if tag in ("insert", "replace"):
            for j in range(j1, j2):
                new_ins[j] = True
    return old_del, new_ins


def _change_spans(old_tokens, new_tokens):
    """Structured before/after spans for display + summaries.
    Each: {"type": "changed"|"added"|"removed", "old": str, "new": str}.
    Boundary-only (whitespace/sentinel) differences are dropped."""
    a = [t["text"] for t in old_tokens]
    b = [t["text"] for t in new_tokens]

    def clean(toks, lo, hi):
        return " ".join(t.strip() for t in toks[lo:hi] if t != SEP and t.strip())

    spans = []
    for tag, i1, i2, j1, j2 in SequenceMatcher(a=a, b=b, autojunk=False).get_opcodes():
        if tag == "equal":
            continue
        old_s, new_s = clean(a, i1, i2), clean(b, j1, j2)
        if not old_s and not new_s:
            continue
        if old_s and new_s:
            spans.append({"type": "changed", "old": old_s, "new": new_s})
        elif new_s:
            spans.append({"type": "added", "old": "", "new": new_s})
        else:
            spans.append({"type": "removed", "old": old_s, "new": ""})
    return spans


# --------------------------------------------------------------------------- #
# Recolor only the changed runs, in place
# --------------------------------------------------------------------------- #

def _segments(text, flags):
    """Split text into maximal (substring, colored) spans of equal flag."""
    out, i, n = [], 0, len(text)
    while i < n:
        j, cur = i, flags[i]
        while j < n and flags[j] == cur:
            j += 1
        out.append((text[i:j], cur))
        i = j
    return out


def _set_run_text(elem, text):
    run = _Run(elem, None)
    run.text = text
    t_el = elem.find(A_T)
    if t_el is not None:
        t_el.set(XML_SPACE, "preserve")  # keep edge whitespace significant


def _split_and_color_run(p, elem, segments, color):
    """Replace elem with one run per segment (cloning rPr so formatting is kept);
    color only the flagged segments."""
    if len(segments) == 1:
        _, colored = segments[0]
        if colored:
            _Run(elem, None).font.color.rgb = color
        return
    idx = list(p).index(elem)
    p.remove(elem)
    for offset, (text, colored) in enumerate(segments):
        new_elem = copy.deepcopy(elem)  # carries rPr -> full formatting preserved
        p.insert(idx + offset, new_elem)
        _set_run_text(new_elem, text)
        if colored:
            _Run(new_elem, None).font.color.rgb = color


def _color_paragraph(info: _ParaInfo, char_flags, color):
    p = info.para._p
    full = info.text
    for (start, end, elem, tag) in info.run_map:
        if tag != A_R:
            continue  # never recolor field/break elements
        seg_flags = char_flags[start:end]
        if not any(seg_flags):
            continue
        _split_and_color_run(p, elem, _segments(full[start:end], seg_flags), color)


def _recolor_in_place(para_infos, tokens, flags, color) -> bool:
    """Paint flagged, non-whitespace tokens. Returns True iff anything was painted
    (this — not raw text inequality — is what drives has_changes)."""
    char_flags_by_p: dict[int, list[bool]] = {}
    for tok, fl in zip(tokens, flags):
        pidx = tok["pidx"]
        if pidx is None:
            continue
        colored = bool(fl) and not tok["ws"]  # whitespace is never colored
        char_flags_by_p.setdefault(pidx, []).extend([colored] * len(tok["text"]))

    painted = False
    for pidx, cflags in char_flags_by_p.items():
        if any(cflags):
            _color_paragraph(para_infos[pidx], cflags, color)
            painted = True
    return painted


# --------------------------------------------------------------------------- #
# Speaker notes (text only; not boxed on the slide image)
# --------------------------------------------------------------------------- #

def _notes_text(slide) -> str:
    if slide is None:
        return ""
    try:
        if slide.has_notes_slide:
            return (slide.notes_slide.notes_text_frame.text or "").strip()
    except Exception:
        pass
    return ""


# --------------------------------------------------------------------------- #
# Optimal slide matching (Hungarian over title/body similarity)
# --------------------------------------------------------------------------- #

def _norm(s: str) -> str:
    return re.sub(r"\s+", " ", (s or "")).strip().lower()


def _slide_text_parts(slide):
    """(title, body, full) normalized — for matching only, never for coloring."""
    if slide is None:
        return "", "", ""
    title, title_txbody = "", None
    try:
        t = slide.shapes.title
        if t is not None and t.has_text_frame:
            title_txbody = t.text_frame._txBody
            title = t.text_frame.text
    except Exception:
        pass
    body = []
    for tf in _iter_text_frames(slide.shapes):
        if title_txbody is not None and tf._txBody is title_txbody:
            continue
        body.append(tf.text)
    return _norm(title), _norm(" ".join(body)), _norm(" ".join([title, *body]))


def _combined_sim(o, n) -> float:
    body = SequenceMatcher(a=o[1], b=n[1], autojunk=False).ratio()
    if o[0] or n[0]:
        title = SequenceMatcher(a=o[0], b=n[0], autojunk=False).ratio()
        return 0.4 * title + 0.6 * body
    return body


def _lis_mask(seq):
    """Boolean mask: True where the position is part of a longest increasing
    subsequence (used to flag out-of-order, i.e. reordered, matches)."""
    n = len(seq)
    if n == 0:
        return []
    tails, tails_idx, prev = [], [], [-1] * n
    for i, x in enumerate(seq):
        pos = bisect.bisect_left(tails, x)
        if pos == len(tails):
            tails.append(x)
            tails_idx.append(i)
        else:
            tails[pos] = x
            tails_idx[pos] = i
        prev[i] = tails_idx[pos - 1] if pos > 0 else -1
    mask = [False] * n
    k = tails_idx[-1]
    while k != -1:
        mask[k] = True
        k = prev[k]
    return mask


def match_slides(old_parts, new_parts, threshold: float = 0.55):
    """Return ordered pairs: each {"old", "new", "kind", "reordered"}.
    kind in {"matched", "added", "removed"}; old/new are 0-based indices or None.
    """
    no, nn = len(old_parts), len(new_parts)
    if no == 0 and nn == 0:
        return []
    if no == 0:
        return [{"old": None, "new": j, "kind": "added", "reordered": False} for j in range(nn)]
    if nn == 0:
        return [{"old": i, "new": None, "kind": "removed", "reordered": False} for i in range(no)]

    sim = np.zeros((no, nn))
    for i, o in enumerate(old_parts):
        for j, n in enumerate(new_parts):
            sim[i, j] = _combined_sim(o, n)
    row_ind, col_ind = linear_sum_assignment(1.0 - sim)  # Hungarian; handles unequal counts

    accepted = sorted(
        [(int(r), int(c)) for r, c in zip(row_ind, col_ind) if sim[r, c] >= threshold],
        key=lambda rc: rc[0],
    )
    order_ok = _lis_mask([c for _, c in accepted])
    matched_old = {r for r, _ in accepted}
    matched_new = {c for _, c in accepted}

    pairs = [
        {"old": r, "new": c, "kind": "matched", "reordered": not order_ok[k]}
        for k, (r, c) in enumerate(accepted)
    ]
    pairs += [
        {"old": i, "new": None, "kind": "removed", "reordered": False}
        for i in range(no) if i not in matched_old
    ]
    pairs += [
        {"old": None, "new": j, "kind": "added", "reordered": False}
        for j in range(nn) if j not in matched_new
    ]
    pairs.sort(key=lambda p: (
        p["new"] if p["new"] is not None else 1e9,
        p["old"] if p["old"] is not None else 1e9,
    ))
    return pairs


# --------------------------------------------------------------------------- #
# Renderer: pptx -> pdf -> png per slide
# --------------------------------------------------------------------------- #

class RenderError(RuntimeError):
    pass


def pptx_to_images(pptx_path, out_dir, dpi: int = 150, timeout: int = 120):
    """Rasterize a deck to one PNG per slide. A unique LibreOffice user profile
    per call keeps concurrent (multi-replica) invocations from deadlocking on the
    shared profile lock."""
    pptx_path = Path(pptx_path)
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    profile = f"file:///tmp/lo_{uuid.uuid4().hex}"

    try:
        subprocess.run(
            ["soffice", f"-env:UserInstallation={profile}", "--headless",
             "--convert-to", "pdf", "--outdir", str(out_dir), str(pptx_path)],
            check=True, capture_output=True, timeout=timeout,
        )
    except subprocess.TimeoutExpired as e:
        raise RenderError(f"LibreOffice timed out after {timeout}s") from e
    except subprocess.CalledProcessError as e:
        raise RenderError(f"LibreOffice failed: {e.stderr.decode('utf-8', 'ignore')[:500]}") from e

    pdf = out_dir / f"{pptx_path.stem}.pdf"
    if not pdf.exists():
        raise RenderError(f"Expected PDF not produced: {pdf}")

    try:
        subprocess.run(
            ["pdftoppm", "-png", "-r", str(dpi), str(pdf), str(out_dir / "slide")],
            check=True, capture_output=True, timeout=timeout,
        )
    except subprocess.TimeoutExpired as e:
        raise RenderError(f"pdftoppm timed out after {timeout}s") from e
    except subprocess.CalledProcessError as e:
        raise RenderError(f"pdftoppm failed: {e.stderr.decode('utf-8', 'ignore')[:500]}") from e

    return sorted(out_dir.glob("slide-*.png"))  # pdftoppm zero-pads -> lexical sort is correct


# --------------------------------------------------------------------------- #
# Orchestration
# --------------------------------------------------------------------------- #

@dataclass
class SlideEntry:
    old_slide: Optional[int]      # 1-based index in old deck, or None
    new_slide: Optional[int]      # 1-based index in new deck, or None
    status: str                   # "changed" | "unchanged" | "added" | "removed"
    possibly_reordered: bool
    has_changes: bool
    has_notes_changes: bool
    text_similarity: float
    changes: list
    old_text: str
    new_text: str
    old_notes: str
    new_notes: str
    old_index0: Optional[int]     # 0-based, internal (for image lookup)
    new_index0: Optional[int]

    def public_dict(self):
        return {
            "old_slide": self.old_slide,
            "new_slide": self.new_slide,
            "status": self.status,
            "possibly_reordered": self.possibly_reordered,
            "has_changes": self.has_changes,
            "has_notes_changes": self.has_notes_changes,
            "text_similarity": round(self.text_similarity, 3),
            "changes": self.changes,
            "old_text": self.old_text,
            "new_text": self.new_text,
            "old_notes": self.old_notes,
            "new_notes": self.new_notes,
        }


def build_marked_decks(old_path, new_path, threshold: float = 0.55):
    """Color both decks in place and compute per-pair metadata.
    Returns (old_prs, new_prs, entries). No rendering happens here, so this is
    fully unit-testable without LibreOffice."""
    old_prs = Presentation(str(old_path))
    new_prs = Presentation(str(new_path))
    old_slides = list(old_prs.slides)
    new_slides = list(new_prs.slides)

    pairs = match_slides(
        [_slide_text_parts(s) for s in old_slides],
        [_slide_text_parts(s) for s in new_slides],
        threshold=threshold,
    )

    entries: list[SlideEntry] = []
    for p in pairs:
        oi, ni = p["old"], p["new"]
        old_slide = old_slides[oi] if oi is not None else None
        new_slide = new_slides[ni] if ni is not None else None

        op, ot = _build_stream(old_slide)
        npar, nt = _build_stream(new_slide)
        old_del, new_ins = _flag(ot, nt)

        painted_old = _recolor_in_place(op, ot, old_del, RED) if old_slide is not None else False
        painted_new = _recolor_in_place(npar, nt, new_ins, GREEN) if new_slide is not None else False

        old_text = _stream_plaintext(ot)
        new_text = _stream_plaintext(nt)
        spans = _change_spans(ot, nt)
        old_notes = _notes_text(old_slide)
        new_notes = _notes_text(new_slide)
        has_notes_changes = _norm(old_notes) != _norm(new_notes)

        if p["kind"] == "added":
            status, has_changes = "added", True
        elif p["kind"] == "removed":
            status, has_changes = "removed", True
        else:
            has_changes = painted_old or painted_new
            status = "changed" if has_changes else "unchanged"

        sim = SequenceMatcher(a=_norm(old_text), b=_norm(new_text), autojunk=False).ratio()

        entries.append(SlideEntry(
            old_slide=(oi + 1) if oi is not None else None,
            new_slide=(ni + 1) if ni is not None else None,
            status=status,
            possibly_reordered=p["reordered"],
            has_changes=has_changes,
            has_notes_changes=has_notes_changes,
            text_similarity=sim,
            changes=spans,
            old_text=old_text,
            new_text=new_text,
            old_notes=old_notes,
            new_notes=new_notes,
            old_index0=oi,
            new_index0=ni,
        ))
    return old_prs, new_prs, entries


def make_diff_previews(old_path, new_path, work_dir, dpi: int = 150, timeout: int = 120):
    """Full pipeline: mark both decks, render each once, attach PNG paths.
    Returns (entries, old_png_paths, new_png_paths)."""
    work = Path(work_dir)
    work.mkdir(parents=True, exist_ok=True)

    old_prs, new_prs, entries = build_marked_decks(old_path, new_path)

    old_marked = work / "old_marked.pptx"
    new_marked = work / "new_marked.pptx"
    old_prs.save(str(old_marked))
    new_prs.save(str(new_marked))

    old_imgs = pptx_to_images(old_marked, work / "old", dpi=dpi, timeout=timeout)
    new_imgs = pptx_to_images(new_marked, work / "new", dpi=dpi, timeout=timeout)
    return entries, old_imgs, new_imgs
